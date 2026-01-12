"""
Report Generation Tools - 报告生成工具集

为 Reporter Agent 提供专业的报告生成能力，包括：
- chart_generator: 生成 Mermaid/ASCII 图表
- table_formatter: 格式化 Markdown 表格
- markdown_exporter: 导出格式化 Markdown 文档
- pdf_exporter: 导出 PDF 文档
- ppt_generator: 生成 PowerPoint 演示文稿

所有工具返回字符串类型，供 LLM Agent 使用。
工具优先生成文本格式输出（Mermaid, Markdown），复杂格式（PDF, PPTX）保存到文件。
"""

import json
from AICrews.observability.logging import get_logger
import os
from datetime import datetime
from typing import Dict, Any, List, Optional, Union

from crewai.tools import tool

logger = get_logger(__name__)

# 配置输出目录
REPORTS_OUTPUT_DIR = os.environ.get("REPORTS_OUTPUT_DIR", "results/reports")


def _ensure_output_dir(subdir: str = "") -> str:
    """确保输出目录存在"""
    path = os.path.join(REPORTS_OUTPUT_DIR, subdir) if subdir else REPORTS_OUTPUT_DIR
    os.makedirs(path, exist_ok=True)
    return path


def _generate_filename(prefix: str, ext: str, ticker: str = "") -> str:
    """生成带时间戳的文件名"""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    ticker_part = f"_{ticker}" if ticker else ""
    return f"{prefix}{ticker_part}_{timestamp}.{ext}"


# =============================================================================
# P0: 基础工具 - chart_generator
# =============================================================================

@tool("chart_generator")
def chart_generator(
    data: str,
    chart_type: str = "line",
    title: str = "",
) -> str:
    """Generate a chart from financial data in Mermaid or ASCII format.

    Creates visual representations of financial data that can be embedded in reports.
    Mermaid charts can be rendered by most Markdown viewers and documentation systems.

    Args:
        data: JSON string with chart data. Formats supported:
              - Time series: {"labels": ["Q1", "Q2", ...], "values": [100, 120, ...]}
              - Multi-series: {"labels": [...], "series": [{"name": "Revenue", "values": [...]}, ...]}
              - XY data: {"x": [...], "y": [...]}
        chart_type: Type of chart - "line", "bar", "pie", "area" (default: "line")
        title: Chart title (optional)

    Returns:
        Mermaid chart syntax or ASCII chart that can be embedded in Markdown

    Example:
        >>> chart_generator('{"labels": ["Q1","Q2","Q3","Q4"], "values": [100,150,120,180]}', "bar", "Quarterly Revenue")
    """
    try:
        # 解析输入数据
        if isinstance(data, str):
            chart_data = json.loads(data)
        else:
            chart_data = data

        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        series = chart_data.get("series", [])

        # 如果是 XY 格式，转换为 labels/values
        if "x" in chart_data and "y" in chart_data:
            labels = chart_data["x"]
            values = chart_data["y"]

        # 如果没有 series，使用单系列
        if not series and values:
            series = [{"name": "Value", "values": values}]

        chart_type_lower = chart_type.lower()

        # 生成 Mermaid 图表
        if chart_type_lower in ["line", "area"]:
            return _generate_mermaid_xy_chart(labels, series, title, chart_type_lower)
        elif chart_type_lower == "bar":
            return _generate_mermaid_bar_chart(labels, series, title)
        elif chart_type_lower == "pie":
            return _generate_mermaid_pie_chart(labels, values, title)
        else:
            # 未知类型，使用 ASCII 表格
            return _generate_ascii_chart(labels, series, title)

    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON data for chart: {e}")
        return f"Error: Invalid JSON data - {str(e)}"
    except Exception as e:
        logger.error(f"Error generating chart: {e}")
        return f"Error generating chart: {str(e)}"


def _generate_mermaid_xy_chart(
    labels: List[str],
    series: List[Dict[str, Any]],
    title: str,
    chart_type: str = "line"
) -> str:
    """生成 Mermaid XY 图表（折线图/面积图）"""
    lines = ["```mermaid", "xychart-beta"]

    if title:
        lines.append(f'    title "{title}"')

    # X 轴标签
    if labels:
        labels_str = ", ".join(f'"{l}"' for l in labels)
        lines.append(f"    x-axis [{labels_str}]")

    # Y 轴（自动推断范围）
    all_values = []
    for s in series:
        all_values.extend(s.get("values", []))

    if all_values:
        min_val = min(all_values)
        max_val = max(all_values)
        padding = (max_val - min_val) * 0.1 or 10
        lines.append(f"    y-axis \"Value\" {min_val - padding:.0f} --> {max_val + padding:.0f}")

    # 数据系列
    line_type = "line" if chart_type == "line" else "area"
    for s in series:
        name = s.get("name", "Data")
        vals = s.get("values", [])
        vals_str = ", ".join(str(v) for v in vals)
        lines.append(f'    {line_type} [{vals_str}]')

    lines.append("```")
    return "\n".join(lines)


def _generate_mermaid_bar_chart(
    labels: List[str],
    series: List[Dict[str, Any]],
    title: str
) -> str:
    """生成 Mermaid 条形图"""
    lines = ["```mermaid", "xychart-beta"]

    if title:
        lines.append(f'    title "{title}"')

    if labels:
        labels_str = ", ".join(f'"{l}"' for l in labels)
        lines.append(f"    x-axis [{labels_str}]")

    # Y 轴
    all_values = []
    for s in series:
        all_values.extend(s.get("values", []))

    if all_values:
        min_val = min(0, min(all_values))  # 条形图通常从 0 开始
        max_val = max(all_values)
        padding = max_val * 0.1 or 10
        lines.append(f"    y-axis \"Value\" {min_val:.0f} --> {max_val + padding:.0f}")

    for s in series:
        vals = s.get("values", [])
        vals_str = ", ".join(str(v) for v in vals)
        lines.append(f'    bar [{vals_str}]')

    lines.append("```")
    return "\n".join(lines)


def _generate_mermaid_pie_chart(
    labels: List[str],
    values: List[float],
    title: str
) -> str:
    """生成 Mermaid 饼图"""
    lines = ["```mermaid", "pie showData"]

    if title:
        lines.append(f'    title {title}')

    for label, value in zip(labels, values):
        lines.append(f'    "{label}" : {value}')

    lines.append("```")
    return "\n".join(lines)


def _generate_ascii_chart(
    labels: List[str],
    series: List[Dict[str, Any]],
    title: str
) -> str:
    """生成 ASCII 文本图表（后备方案）"""
    lines = []
    if title:
        lines.append(f"=== {title} ===")
        lines.append("")

    # 简单的 ASCII 条形图
    for s in series:
        name = s.get("name", "Data")
        vals = s.get("values", [])

        if not vals:
            continue

        max_val = max(vals) if vals else 1

        lines.append(f"[{name}]")
        for i, (label, val) in enumerate(zip(labels, vals)):
            bar_len = int((val / max_val) * 40) if max_val else 0
            bar = "█" * bar_len
            lines.append(f"  {label:10s} | {bar} {val}")
        lines.append("")

    return "\n".join(lines)


# =============================================================================
# P0: 基础工具 - table_formatter
# =============================================================================

@tool("table_formatter")
def table_formatter(
    data: str,
    columns: str = "",
    title: str = "",
    format_rules: str = "",
) -> str:
    """Format financial data as a professional Markdown table.

    Creates well-formatted tables with proper alignment and optional value formatting
    for currencies, percentages, and numbers.

    Args:
        data: JSON string with table data as list of row dictionaries.
              Example: '[{"ticker": "AAPL", "price": 150.25, "change": 0.025}]'
        columns: Comma-separated column names to include (default: all columns).
                 Example: "ticker,price,change"
        title: Optional table title
        format_rules: JSON string with formatting rules.
                     Example: '{"price": "currency", "change": "percent"}'
                     Supported formats: "currency", "percent", "number", "integer"

    Returns:
        Formatted Markdown table string

    Example:
        >>> table_formatter('[{"ticker":"AAPL","price":150.25}]', "ticker,price", "Stock Prices", '{"price":"currency"}')
    """
    try:
        # 解析输入
        rows = json.loads(data) if isinstance(data, str) else data

        if not rows:
            return "_No data available_"

        # 解析列
        if columns:
            col_list = [c.strip() for c in columns.split(",")]
        else:
            # 使用第一行的所有键
            col_list = list(rows[0].keys())

        # 解析格式规则
        rules = {}
        if format_rules:
            rules = json.loads(format_rules) if isinstance(format_rules, str) else format_rules

        # 构建表格
        lines = []

        if title:
            lines.append(f"**{title}**")
            lines.append("")

        # 表头
        header = "| " + " | ".join(col_list) + " |"
        lines.append(header)

        # 分隔线（根据数据类型设置对齐）
        alignments = []
        for col in col_list:
            fmt = rules.get(col, "")
            if fmt in ["currency", "percent", "number", "integer"]:
                alignments.append(":---:")  # 数字居中
            else:
                alignments.append(":---")   # 文本左对齐

        separator = "| " + " | ".join(alignments) + " |"
        lines.append(separator)

        # 数据行
        for row in rows:
            cells = []
            for col in col_list:
                value = row.get(col, "")
                formatted = _format_cell_value(value, rules.get(col, ""))
                cells.append(formatted)
            lines.append("| " + " | ".join(cells) + " |")

        return "\n".join(lines)

    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON data for table: {e}")
        return f"Error: Invalid JSON data - {str(e)}"
    except Exception as e:
        logger.error(f"Error formatting table: {e}")
        return f"Error formatting table: {str(e)}"


def _format_cell_value(value: Any, format_type: str) -> str:
    """根据格式规则格式化单元格值"""
    if value is None or value == "":
        return "-"

    try:
        if format_type == "currency":
            num = float(value)
            if num >= 0:
                return f"${num:,.2f}"
            else:
                return f"-${abs(num):,.2f}"

        elif format_type == "percent":
            num = float(value)
            # 假设输入是小数形式 (0.05 = 5%)
            if abs(num) < 1:
                num = num * 100
            sign = "+" if num > 0 else ""
            return f"{sign}{num:.2f}%"

        elif format_type == "number":
            num = float(value)
            return f"{num:,.2f}"

        elif format_type == "integer":
            num = int(float(value))
            return f"{num:,}"

        else:
            return str(value)

    except (ValueError, TypeError):
        return str(value)


# =============================================================================
# P0: 基础工具 - markdown_exporter
# =============================================================================

@tool("markdown_exporter")
def markdown_exporter(
    sections: str,
    template: str = "standard",
    metadata: str = "",
) -> str:
    """Export report sections to a formatted Markdown document.

    Assembles multiple sections into a cohesive, professionally formatted report
    with proper headings, metadata, and structure.

    Args:
        sections: JSON string with section definitions.
                  Format: '[{"title": "Summary", "content": "...", "level": 1}]'
                  - title: Section title
                  - content: Section content (can include Markdown)
                  - level: Heading level (1-4, default: 2)
        template: Document template style
                  - "standard": Basic report structure
                  - "executive": Executive summary format with highlights
                  - "detailed": Comprehensive format with TOC
        metadata: JSON string with document metadata.
                  Example: '{"ticker": "AAPL", "date": "2024-01-15", "author": "AI Analyst"}'

    Returns:
        Complete Markdown document as string

    Example:
        >>> markdown_exporter('[{"title":"Overview","content":"Analysis of AAPL","level":1}]', "standard", '{"ticker":"AAPL"}')
    """
    try:
        # 解析输入
        section_list = json.loads(sections) if isinstance(sections, str) else sections
        meta = json.loads(metadata) if metadata and isinstance(metadata, str) else (metadata or {})

        lines = []

        # 文档头部
        if template == "executive":
            lines.extend(_generate_executive_header(meta))
        elif template == "detailed":
            lines.extend(_generate_detailed_header(meta, section_list))
        else:
            lines.extend(_generate_standard_header(meta))

        # 内容部分
        for section in section_list:
            title = section.get("title", "")
            content = section.get("content", "")
            level = section.get("level", 2)

            # 限制标题级别
            level = max(1, min(4, level))
            heading = "#" * level

            lines.append(f"{heading} {title}")
            lines.append("")
            lines.append(content)
            lines.append("")

        # 文档脚注
        lines.extend(_generate_footer(meta, template))

        return "\n".join(lines)

    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON for markdown export: {e}")
        return f"Error: Invalid JSON - {str(e)}"
    except Exception as e:
        logger.error(f"Error exporting markdown: {e}")
        return f"Error exporting markdown: {str(e)}"


def _generate_standard_header(meta: Dict[str, Any]) -> List[str]:
    """生成标准报告头部"""
    lines = []

    ticker = meta.get("ticker", "")
    date = meta.get("date", datetime.now().strftime("%Y-%m-%d"))

    if ticker:
        lines.append(f"# {ticker} Analysis Report")
    else:
        lines.append("# Financial Analysis Report")

    lines.append("")
    lines.append(f"**Date**: {date}")

    if meta.get("author"):
        lines.append(f"**Author**: {meta['author']}")

    lines.append("")
    lines.append("---")
    lines.append("")

    return lines


def _generate_executive_header(meta: Dict[str, Any]) -> List[str]:
    """生成执行摘要风格头部"""
    lines = []

    ticker = meta.get("ticker", "")
    date = meta.get("date", datetime.now().strftime("%Y-%m-%d"))

    lines.append("---")
    lines.append(f"title: {ticker or 'Financial'} Executive Summary")
    lines.append(f"date: {date}")
    if meta.get("author"):
        lines.append(f"author: {meta['author']}")
    lines.append("---")
    lines.append("")

    return lines


def _generate_detailed_header(meta: Dict[str, Any], sections: List[Dict]) -> List[str]:
    """生成详细报告头部（含目录）"""
    lines = _generate_standard_header(meta)

    # 添加目录
    lines.append("## Table of Contents")
    lines.append("")

    for i, section in enumerate(sections, 1):
        title = section.get("title", f"Section {i}")
        # 生成 anchor
        anchor = title.lower().replace(" ", "-").replace(".", "")
        lines.append(f"{i}. [{title}](#{anchor})")

    lines.append("")
    lines.append("---")
    lines.append("")

    return lines


def _generate_footer(meta: Dict[str, Any], template: str) -> List[str]:
    """生成文档脚注"""
    lines = []
    lines.append("---")
    lines.append("")

    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    lines.append(f"*Generated on {timestamp}*")

    if template == "executive":
        lines.append("")
        lines.append("*This report is for informational purposes only and does not constitute financial advice.*")

    return lines


# =============================================================================
# P1: PDF 导出
# =============================================================================

@tool("pdf_exporter")
def pdf_exporter(
    markdown_content: str,
    output_filename: str = "",
    style: str = "professional",
) -> str:
    """Export Markdown content to a PDF document.

    Converts Markdown report to PDF with professional styling.
    Requires weasyprint or similar PDF library to be installed.

    Args:
        markdown_content: Markdown formatted report content
        output_filename: Output PDF filename (auto-generated if empty)
        style: PDF style template
               - "professional": Clean, corporate style
               - "minimal": Simple, minimal formatting
               - "detailed": Detailed with more structure

    Returns:
        Path to generated PDF file, or error message if PDF generation fails.
        Falls back to returning styled HTML if PDF libraries unavailable.

    Example:
        >>> pdf_exporter("# Report\\n\\nContent here...", "analysis_report.pdf", "professional")
    """
    try:
        # 尝试导入 PDF 生成库
        try:
            import markdown
            HAS_MARKDOWN = True
        except ImportError:
            HAS_MARKDOWN = False

        try:
            from weasyprint import HTML, CSS
            HAS_WEASYPRINT = True
        except ImportError:
            HAS_WEASYPRINT = False

        # 先将 Markdown 转换为 HTML
        if HAS_MARKDOWN:
            html_content = markdown.markdown(
                markdown_content,
                extensions=['tables', 'fenced_code', 'toc']
            )
        else:
            # 简单的 Markdown 到 HTML 转换
            html_content = _simple_md_to_html(markdown_content)

        # 获取样式
        css_content = _get_pdf_style(style)

        # 完整 HTML 文档
        full_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>{css_content}</style>
</head>
<body>
    <div class="container">
        {html_content}
    </div>
</body>
</html>"""

        if not HAS_WEASYPRINT:
            # 如果没有 weasyprint，返回 HTML 并提示
            logger.warning("weasyprint not installed. Returning HTML content.")

            # 保存 HTML 文件作为替代
            output_dir = _ensure_output_dir("html")
            if not output_filename:
                output_filename = _generate_filename("report", "html")
            else:
                output_filename = output_filename.replace(".pdf", ".html")

            output_path = os.path.join(output_dir, output_filename)

            with open(output_path, "w", encoding="utf-8") as f:
                f.write(full_html)

            return f"PDF export requires weasyprint library. HTML saved to: {output_path}"

        # 生成 PDF
        output_dir = _ensure_output_dir("pdf")
        if not output_filename:
            output_filename = _generate_filename("report", "pdf")

        output_path = os.path.join(output_dir, output_filename)

        html_doc = HTML(string=full_html)
        html_doc.write_pdf(output_path)

        logger.info(f"PDF exported to: {output_path}")
        return f"PDF successfully exported to: {output_path}"

    except Exception as e:
        logger.error(f"Error exporting PDF: {e}")
        return f"Error exporting PDF: {str(e)}"


def _simple_md_to_html(md: str) -> str:
    """简单的 Markdown 到 HTML 转换（无依赖）"""
    import re

    html = md

    # 标题
    html = re.sub(r'^#### (.+)$', r'<h4>\1</h4>', html, flags=re.MULTILINE)
    html = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html, flags=re.MULTILINE)
    html = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html, flags=re.MULTILINE)
    html = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html, flags=re.MULTILINE)

    # 粗体和斜体
    html = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html)
    html = re.sub(r'\*(.+?)\*', r'<em>\1</em>', html)

    # 列表
    html = re.sub(r'^- (.+)$', r'<li>\1</li>', html, flags=re.MULTILINE)
    html = re.sub(r'(<li>.*</li>\n?)+', r'<ul>\g<0></ul>', html)

    # 分隔线
    html = re.sub(r'^---$', '<hr>', html, flags=re.MULTILINE)

    # 段落
    paragraphs = html.split('\n\n')
    html = '\n'.join(
        f'<p>{p}</p>' if not p.startswith('<') else p
        for p in paragraphs if p.strip()
    )

    return html


def _get_pdf_style(style: str) -> str:
    """获取 PDF 样式"""
    base_style = """
        body {
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
            line-height: 1.6;
            color: #333;
            max-width: 800px;
            margin: 0 auto;
            padding: 20px;
        }
        h1 { color: #1a1a1a; border-bottom: 2px solid #333; padding-bottom: 10px; }
        h2 { color: #2a2a2a; border-bottom: 1px solid #ddd; padding-bottom: 5px; }
        h3 { color: #3a3a3a; }
        table { border-collapse: collapse; width: 100%; margin: 20px 0; }
        th, td { border: 1px solid #ddd; padding: 8px 12px; text-align: left; }
        th { background-color: #f5f5f5; font-weight: bold; }
        tr:nth-child(even) { background-color: #fafafa; }
        code { background-color: #f4f4f4; padding: 2px 6px; border-radius: 3px; }
        pre { background-color: #f4f4f4; padding: 15px; border-radius: 5px; overflow-x: auto; }
        blockquote { border-left: 4px solid #ddd; margin: 0; padding-left: 15px; color: #666; }
        hr { border: none; border-top: 1px solid #ddd; margin: 30px 0; }
    """

    if style == "minimal":
        return base_style + """
            body { font-family: Georgia, serif; }
            h1 { border-bottom: none; }
            h2 { border-bottom: none; }
        """
    elif style == "detailed":
        return base_style + """
            .container { padding: 40px; }
            h1 { font-size: 28px; }
            h2 { font-size: 22px; margin-top: 30px; }
            table { font-size: 14px; }
        """
    else:  # professional
        return base_style


# =============================================================================
# P2: PPT 生成
# =============================================================================

@tool("ppt_generator")
def ppt_generator(
    slides: str,
    template: str = "financial",
    output_filename: str = "",
) -> str:
    """Generate a PowerPoint presentation from structured data.

    Creates professional presentations for financial reports and analyses.
    Requires python-pptx library to be installed.

    Args:
        slides: JSON string with slide definitions.
                Format: '[{"title": "Slide Title", "content": "...", "layout": "title_content", "bullets": [...]}]'
                Supported layouts:
                - "title": Title slide
                - "title_content": Title with text content
                - "bullets": Title with bullet points
                - "two_column": Two-column layout
                - "chart": Chart placeholder (renders as text description)
        template: Presentation template
                  - "financial": Finance-focused styling
                  - "minimal": Clean, minimal design
                  - "executive": Executive presentation style
        output_filename: Output PPTX filename (auto-generated if empty)

    Returns:
        Path to generated PPTX file, or text representation if library unavailable

    Example:
        >>> ppt_generator('[{"title":"Q4 Results","content":"Revenue up 15%","layout":"title_content"}]', "financial")
    """
    try:
        # 解析 slides 数据
        slide_list = json.loads(slides) if isinstance(slides, str) else slides

        if not slide_list:
            return "Error: No slides provided"

        # 尝试导入 python-pptx
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN
            HAS_PPTX = True
        except ImportError:
            HAS_PPTX = False

        if not HAS_PPTX:
            # 返回文本表示
            logger.warning("python-pptx not installed. Returning text representation.")
            return _generate_ppt_text_representation(slide_list)

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 比例
        prs.slide_height = Inches(7.5)

        # 添加幻灯片
        for slide_data in slide_list:
            _add_slide_to_presentation(prs, slide_data, template)

        # 保存文件
        output_dir = _ensure_output_dir("pptx")
        if not output_filename:
            output_filename = _generate_filename("presentation", "pptx")

        output_path = os.path.join(output_dir, output_filename)
        prs.save(output_path)

        logger.info(f"PPTX exported to: {output_path}")
        return f"PowerPoint successfully exported to: {output_path}"

    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON for PPT: {e}")
        return f"Error: Invalid JSON - {str(e)}"
    except Exception as e:
        logger.error(f"Error generating PPT: {e}")
        return f"Error generating PPT: {str(e)}"


def _add_slide_to_presentation(prs, slide_data: Dict[str, Any], template: str) -> None:
    """向演示文稿添加幻灯片"""
    from pptx.util import Inches, Pt

    title = slide_data.get("title", "")
    content = slide_data.get("content", "")
    layout = slide_data.get("layout", "title_content")
    bullets = slide_data.get("bullets", [])

    # 选择布局
    if layout == "title":
        slide_layout = prs.slide_layouts[6]  # 空白布局
    else:
        slide_layout = prs.slide_layouts[6]  # 使用空白布局，手动添加内容

    slide = prs.slides.add_slide(slide_layout)

    # 添加标题
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(36 if layout == "title" else 28)
        title_para.font.bold = True

    # 添加内容
    if layout == "title":
        # 仅标题幻灯片，添加副标题
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3), Inches(12.333), Inches(1)
            )
            content_frame = content_box.text_frame
            content_para = content_frame.paragraphs[0]
            content_para.text = content
            content_para.font.size = Pt(24)

    elif layout == "bullets" and bullets:
        # 项目符号列表
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(12.333), Inches(5)
        )
        content_frame = content_box.text_frame

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.space_after = Pt(12)

    elif layout == "two_column":
        # 两列布局
        left_content = slide_data.get("left_content", content)
        right_content = slide_data.get("right_content", "")

        # 左列
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(5.9), Inches(5)
        )
        left_frame = left_box.text_frame
        left_frame.paragraphs[0].text = left_content
        left_frame.paragraphs[0].font.size = Pt(18)

        # 右列
        right_box = slide.shapes.add_textbox(
            Inches(6.9), Inches(2), Inches(5.9), Inches(5)
        )
        right_frame = right_box.text_frame
        right_frame.paragraphs[0].text = right_content
        right_frame.paragraphs[0].font.size = Pt(18)

    else:
        # 默认：标题 + 内容
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(12.333), Inches(5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_para = content_frame.paragraphs[0]
            content_para.text = content
            content_para.font.size = Pt(18)


def _generate_ppt_text_representation(slides: List[Dict[str, Any]]) -> str:
    """生成 PPT 的文本表示（当 python-pptx 不可用时）"""
    lines = ["# Presentation Outline", ""]
    lines.append("*Note: python-pptx library not installed. Text representation below:*")
    lines.append("")

    for i, slide in enumerate(slides, 1):
        title = slide.get("title", f"Slide {i}")
        content = slide.get("content", "")
        layout = slide.get("layout", "title_content")
        bullets = slide.get("bullets", [])

        lines.append(f"## Slide {i}: {title}")
        lines.append(f"*Layout: {layout}*")
        lines.append("")

        if content:
            lines.append(content)
            lines.append("")

        if bullets:
            for bullet in bullets:
                lines.append(f"- {bullet}")
            lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


# =============================================================================
# P3: Infograph 生成 - AntV Infographic DSL
# =============================================================================

@tool("infograph_generator")
def infograph_generator(
    report_data: str,
    sections: str = "all",
    style: str = "professional",
) -> str:
    """Generate AntV Infographic DSL for visual report rendering.

    Creates DSL syntax that can be rendered by @antv/infographic library
    into beautiful visual reports. The DSL is optimized for financial analysis.

    Args:
        report_data: JSON string with report data. Format:
            {
                "ticker": "AAPL",           # Stock ticker (optional)
                "crew_name": "Analysis Crew", # Crew name (optional)
                "metrics": {                 # Key metrics (optional)
                    "Current Price": "$180.50",
                    "P/E Ratio": 28.5,
                    "Market Cap": "2.8T"
                },
                "insights": [                # Key insights list (optional)
                    "Strong revenue growth in Q4",
                    "New product launches driving demand"
                ],
                "recommendation": "buy",     # buy|sell|hold|neutral (optional)
                "stats": {                   # Execution stats (optional)
                    "duration": "2m 30s",
                    "tokens": 15000,
                    "tool_calls": 12,
                    "agents": 3
                }
            }
        sections: Comma-separated sections to include. Options:
            - "all": Include all available sections (default)
            - "header": Title/header section
            - "stats": Execution statistics
            - "metrics": Key metrics cards
            - "insights": Insights list
            - "recommendation": Buy/sell/hold recommendation
            Example: "header,metrics,recommendation"
        style: Visual style theme
            - "professional": Clean corporate style (default)
            - "dark": Dark theme with gradients
            - "minimal": Simple, minimal design

    Returns:
        AntV Infographic DSL string that can be rendered by the frontend

    Example:
        >>> infograph_generator('{"ticker":"AAPL","metrics":{"Price":"$180"},"recommendation":"buy"}')
    """
    try:
        # 解析输入数据
        data = json.loads(report_data) if isinstance(report_data, str) else report_data

        # 解析要包含的部分
        if sections == "all":
            include_sections = ["header", "stats", "metrics", "insights", "recommendation"]
        else:
            include_sections = [s.strip().lower() for s in sections.split(",")]

        # 生成各部分 DSL
        dsl_parts = []

        # Header section
        if "header" in include_sections:
            header_dsl = _generate_infograph_header(
                ticker=data.get("ticker"),
                crew_name=data.get("crew_name"),
                style=style
            )
            if header_dsl:
                dsl_parts.append(header_dsl)

        # Stats section
        if "stats" in include_sections and data.get("stats"):
            stats_dsl = _generate_infograph_stats(data["stats"], style)
            if stats_dsl:
                dsl_parts.append(stats_dsl)

        # Metrics section
        if "metrics" in include_sections and data.get("metrics"):
            metrics_dsl = _generate_infograph_metrics(data["metrics"], style)
            if metrics_dsl:
                dsl_parts.append(metrics_dsl)

        # Insights section
        if "insights" in include_sections and data.get("insights"):
            insights_dsl = _generate_infograph_insights(data["insights"], style)
            if insights_dsl:
                dsl_parts.append(insights_dsl)

        # Recommendation section
        if "recommendation" in include_sections and data.get("recommendation"):
            rec_dsl = _generate_infograph_recommendation(
                recommendation=data["recommendation"],
                ticker=data.get("ticker"),
                style=style
            )
            if rec_dsl:
                dsl_parts.append(rec_dsl)

        if not dsl_parts:
            return "Error: No valid data provided for infograph generation"

        return "\n\n".join(dsl_parts)

    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON for infograph: {e}")
        return f"Error: Invalid JSON - {str(e)}"
    except Exception as e:
        logger.error(f"Error generating infograph: {e}")
        return f"Error generating infograph: {str(e)}"


def _generate_infograph_header(
    ticker: Optional[str],
    crew_name: Optional[str],
    style: str
) -> str:
    """生成 Infograph 头部 DSL"""
    title = f"{ticker} Analysis Report" if ticker else "Analysis Report"
    subtitle = crew_name or "AI-Powered Financial Analysis"
    timestamp = datetime.now().strftime("%B %d, %Y")

    # 根据风格选择背景
    bg_style = {
        "professional": "gradient-light",
        "dark": "gradient-dark",
        "minimal": "solid-white"
    }.get(style, "gradient-light")

    return f"""infographic header-hero
data
  title {title}
  subtitle {subtitle}
  timestamp {timestamp}
style
  background {bg_style}
  accent green"""


def _generate_infograph_stats(stats: Dict[str, Any], style: str) -> str:
    """生成执行统计 DSL"""
    items = []

    stat_configs = [
        ("duration", "Duration", "clock"),
        ("tokens", "Tokens", "zap"),
        ("tool_calls", "Tool Calls", "wrench"),
        ("agents", "Agents", "users"),
    ]

    for key, label, icon in stat_configs:
        if key in stats:
            value = stats[key]
            # 格式化数值
            if isinstance(value, (int, float)) and key == "tokens":
                value = f"{value:,}"
            items.append(f"    - label {label}\n      value {value}\n      icon {icon}")

    if not items:
        return ""

    card_style = "glass" if style == "dark" else "card"

    return f"""infographic stat-card-group
data
  items
{chr(10).join(items)}
style
  layout horizontal
  cardStyle {card_style}"""


def _generate_infograph_metrics(metrics: Dict[str, Any], style: str) -> str:
    """生成指标卡片 DSL"""
    items = []

    for label, value in list(metrics.items())[:6]:  # 最多 6 个指标
        items.append(f"    - label {label}\n      value {value}")

    if not items:
        return ""

    return f"""infographic metric-cards
data
  items
{chr(10).join(items)}
style
  layout grid
  columns 3
  highlight first"""


def _generate_infograph_insights(insights: List[str], style: str) -> str:
    """生成洞察列表 DSL"""
    items = []

    for i, insight in enumerate(insights[:5], 1):  # 最多 5 条洞察
        # 截断过长的洞察
        truncated = insight[:150] if len(insight) > 150 else insight
        items.append(f"    - label Insight {i}\n      desc {truncated}")

    if not items:
        return ""

    accent = "yellow" if style != "minimal" else "gray"

    return f"""infographic list-row-simple-vertical
data
  title Key Insights
  items
{chr(10).join(items)}
style
  icon lightbulb
  accent {accent}"""


def _generate_infograph_recommendation(
    recommendation: str,
    ticker: Optional[str],
    style: str
) -> str:
    """生成投资建议 DSL"""
    rec_lower = recommendation.lower()

    labels = {
        "buy": "BUY",
        "sell": "SELL",
        "hold": "HOLD",
        "neutral": "NEUTRAL",
    }

    colors = {
        "buy": "green",
        "sell": "red",
        "hold": "blue",
        "neutral": "gray",
    }

    descriptions = {
        "buy": "Analysis indicates positive momentum and growth potential",
        "sell": "Analysis suggests reducing position or exiting",
        "hold": "Maintain current position while monitoring developments",
        "neutral": "Insufficient signals for directional bias",
    }

    label = labels.get(rec_lower, rec_lower.upper())
    color = colors.get(rec_lower, "gray")
    desc = descriptions.get(rec_lower, "See detailed analysis for rationale")
    subtitle = ticker or "Asset"

    return f"""infographic callout-highlight
data
  title Recommendation
  value {label}
  subtitle {subtitle}
  description {desc}
style
  accent {color}
  size large
  emphasis strong"""


# =============================================================================
# 工具导出
# =============================================================================

__all__ = [
    "chart_generator",
    "table_formatter",
    "markdown_exporter",
    "pdf_exporter",
    "ppt_generator",
    "infograph_generator",
]
